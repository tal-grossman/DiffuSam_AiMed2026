#!/usr/bin/env python3
"""Generate AIMed 2026 DiffuSAM poster as PPTX (A0 portrait)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

PROJ = os.path.dirname(os.path.abspath(__file__))

# A0 portrait in EMU (841 mm × 1189 mm)
A0_W = Cm(84.1)
A0_H = Cm(118.9)

# Colours (matching LaTeX poster)
DARK_BLUE = RGBColor(0x0E, 0x28, 0x41)
SEC_BLUE = RGBColor(0x15, 0x60, 0x82)
ACCENT_ORG = RGBColor(0xE9, 0x71, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)

MARGIN_L = Cm(2.0)
MARGIN_R = Cm(2.0)
COL_GAP = Cm(1.4)
BODY_W = A0_W - MARGIN_L - MARGIN_R
COL_W = (BODY_W - COL_GAP) // 2

prs = Presentation()
prs.slide_width = A0_W
prs.slide_height = A0_H

slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)


def add_rect(left, top, width, height, fill_color=None, line_color=None, line_w=None):
    from pptx.util import Emu as _E
    shape = slide.shapes.add_shape(
        1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w or Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_para(para, text, size=Pt(27), bold=False, italic=False,
             color=BLACK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    para.alignment = alignment
    run = para.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return run


def section_bar(left, top, width, text):
    """Blue rounded-rectangle section header."""
    shape = slide.shapes.add_shape(
        5, left, top, width, Cm(1.6))  # ROUNDED_RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = SEC_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(35)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = 'Calibri'
    return Cm(1.6) + Cm(0.3)


def body_text(left, top, width, paragraphs, size=Pt(24)):
    """Add a textbox with multiple paragraphs. Returns height used."""
    est_h = Cm(len(paragraphs) * 1.5 + 1)
    tb = add_textbox(left, top, width, est_h)
    tf = tb.text_frame
    tf.word_wrap = True

    for i, (txt, bold, italic) in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = txt
        run.font.size = size
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = BLACK
        run.font.name = 'Calibri'

    tb.text_frame.auto_size = None
    return tb


def add_image(left, top, width, filename):
    path = os.path.join(PROJ, filename)
    if os.path.exists(path):
        pic = slide.shapes.add_picture(path, left, top, width=width)
        return pic.height
    return Cm(3)


# ================================================================
#  HEADER BANNER
# ================================================================
HEADER_H = Cm(10.5)
add_rect(0, 0, A0_W, HEADER_H, fill_color=DARK_BLUE)
add_rect(0, HEADER_H, A0_W, Pt(5), fill_color=ACCENT_ORG)

# Logo left
add_image(Cm(1.5), Cm(1.2), Cm(11), 'tau_logo.png')
# Logo right
add_image(A0_W - Cm(12.5), Cm(1.2), Cm(11), 'tau_logo.png')

# Title
tb = add_textbox(Cm(14), Cm(0.8), Cm(56), Cm(5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "DiffuSAM: Diffusion-Based Prompt-Free SAM2 for\nFew-Shot and Source-Free Medical Image Segmentation"
run.font.size = Pt(60)
run.font.bold = True
run.font.color.rgb = WHITE
run.font.name = 'Calibri'

# Authors
tb = add_textbox(Cm(8), Cm(6.0), Cm(68), Cm(2.0))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Tal Grossman¹†    Noa Cahan¹†    Lev Ayzenberg¹    Hayit Greenspan²"
run.font.size = Pt(38)
run.font.bold = False
run.font.color.rgb = WHITE
run.font.name = 'Calibri'

# Affiliations
tb = add_textbox(Cm(4), Cm(8.0), Cm(76), Cm(2.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = ("†Equal contribution    "
            "¹School of Electrical and Computer Engineering, Tel Aviv University, Israel    "
            "²Dept. of Biomedical Engineering, Tel Aviv University, Israel    "
            "Contact: talgrossman22@gmail.com")
run.font.size = Pt(24)
run.font.color.rgb = WHITE
run.font.name = 'Calibri'

# ================================================================
#  KEYWORDS BAR
# ================================================================
KW_TOP = HEADER_H + Pt(8)
kw_shape = add_rect(MARGIN_L, KW_TOP, BODY_W, Cm(1.8), fill_color=LIGHT_GRAY)
tb = add_textbox(MARGIN_L + Cm(0.5), KW_TOP + Cm(0.2), BODY_W - Cm(1), Cm(1.4))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Keywords: "
run.font.size = Pt(24)
run.font.bold = True
run.font.name = 'Calibri'
run = p.add_run()
run.text = "Medical image segmentation, diffusion models, SAM2, prompt-free segmentation, domain adaptation"
run.font.size = Pt(24)
run.font.name = 'Calibri'

# ================================================================
#  COLUMN POSITIONS
# ================================================================
COL1_L = MARGIN_L
COL2_L = MARGIN_L + COL_W + COL_GAP
BODY_TOP = KW_TOP + Cm(2.5)

# ================================================================
#  LEFT COLUMN
# ================================================================
y = BODY_TOP

# Key Information
dy = section_bar(COL1_L, y, COL_W, "Key Information")
y += dy

txt_key = [
    ("Research question", True, False),
    ("Can a lightweight diffusion prior over frozen SAM2 features enable "
     "accurate, prompt-free medical image segmentation without full "
     "fine-tuning or expert-provided prompts?", False, False),
    ("", False, False),
    ("Findings", True, False),
    ("DiffuSAM achieves 87.2% average Dice in few-shot CT segmentation "
     "and 85.2% in source-free CT-to-MRI domain adaptation on abdominal "
     "organs, competitive with state-of-the-art methods while requiring "
     "<5 GB GPU memory and no user prompts.", False, False),
    ("", False, False),
    ("Meaning", True, False),
    ("Diffusion-based memory generation offers a practical, lightweight "
     "path for adapting vision foundation models to medical imaging "
     "without prompts or expensive fine-tuning, facilitating clinical "
     "workflow integration.", False, False),
]
body_text(COL1_L, y, COL_W, txt_key)
y += Cm(14.5)

# Introduction
dy = section_bar(COL1_L, y, COL_W, "Introduction")
y += dy

txt_intro = [
    ("Medical image segmentation is vital for clinical diagnosis and "
     "treatment planning. Foundation models such as SAM [1] and SAM2 [2] "
     "achieve strong zero-shot performance on natural images but struggle "
     "with medical data due to domain differences in texture, contrast, "
     "and anatomy. Their reliance on expert-provided prompts further "
     "limits clinical deployment. Prior adaptations such as MedSAM [3] "
     "require costly fine-tuning and typically still depend on prompts.",
     False, False),
    ("", False, False),
    ("We propose DiffuSAM, a diffusion-based framework for prompt-free "
     "medical image segmentation that keeps the SAM2 backbone entirely frozen.",
     True, False),
]
body_text(COL1_L, y, COL_W, txt_intro)
y += Cm(10.5)

# Material and Methods
dy = section_bar(COL1_L, y, COL_W, "Material and Methods")
y += dy

txt_methods = [
    ("DiffuSAM trains a lightweight UNet-based diffusion prior conditioned "
     "on frozen SAM2 image encoder features and class labels to generate "
     "SAM2-compatible memory embeddings. These embeddings are injected into "
     "SAM2's memory attention and mask decoder, producing segmentation masks "
     "without user prompts.", False, False),
    ("", False, False),
    ("For volumetric consistency in 3D scans, the prior is additionally "
     "conditioned on predictions from adjacent slices, propagating "
     "bidirectionally from the center slice. The SAM2 backbone remains "
     "entirely frozen; only the compact diffusion prior is trained, "
     "requiring <5 GB GPU memory and converging in under 8,000 iterations.",
     False, False),
]
body_text(COL1_L, y, COL_W, txt_methods)
y += Cm(10.5)

# Pipeline figure (a)
img_h = add_image(COL1_L, y, COL_W, 'pipeline_a_a.png')
y += img_h + Cm(0.3)
tb = add_textbox(COL1_L, y, COL_W, Cm(3.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = ("(a) DiffuSAM overview. A diffusion prior (green) generates memory "
            "embeddings from SAM2 image encoder features, eliminating user "
            "prompts (red). All SAM2 components remain frozen.")
run.font.size = Pt(20)
run.font.italic = True
run.font.name = 'Calibri'
y += Cm(3.5)

# Pipeline figure (b)
img_h = add_image(COL1_L + Cm(3), y, COL_W - Cm(6), 'pipeline_b.png')
y += img_h + Cm(0.3)
tb = add_textbox(COL1_L, y, COL_W, Cm(3.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = ("(b) Diffusion prior model. During training: noised memory embedding, "
            "image embedding, timestep, and class label → learns to denoise. "
            "At inference: generates memory from pure noise in 2 steps.")
run.font.size = Pt(20)
run.font.italic = True
run.font.name = 'Calibri'

# ================================================================
#  RIGHT COLUMN
# ================================================================
y = BODY_TOP

# Results
dy = section_bar(COL2_L, y, COL_W, "Results")
y += dy

txt_results = [
    ("We evaluated DiffuSAM on the BTCV [9] (CT) and CHAOS [10] (MRI) "
     "abdominal datasets for spleen, kidney, and liver segmentation under:",
     False, False),
    ("• Few-shot: 3 of 30 training volumes (BTCV)", False, False),
    ("• Source-free UDA: CT → MRI without target labels", False, False),
]
body_text(COL2_L, y, COL_W, txt_results)
y += Cm(5.5)

# Table 1 — Few-shot
tb = add_textbox(COL2_L, y, COL_W, Cm(1.2))
tf = tb.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Table 1: Few-shot Dice (%) on BTCV (3 training volumes)"
run.font.size = Pt(22)
run.font.bold = True
run.font.name = 'Calibri'
y += Cm(1.4)

table_data_1 = [
    ["Method", "Spleen", "R.Kid.", "L.Kid.", "Liver", "Avg"],
    ["UNet [4]", "54.9", "56.0", "48.7", "85.9", "61.4"],
    ["UNETR [5]", "8.7", "0.2", "2.1", "78.2", "22.3"],
    ["SAM [1]", "15.5", "52.2", "62.2", "26.6", "39.1"],
    ["SAM2 [2]", "20.2", "55.7", "66.9", "34.6", "44.4"],
    ["FATE-SAM [6]", "86.2", "91.0*", "89.3*", "82.2", "87.2"],
    ["DiffuSAM (w/o 3D)", "82.9", "81.0", "83.6", "83.5", "82.7"],
    ["DiffuSAM", "86.5*", "88.6", "86.5", "87.3*", "87.2*"],
]

rows, cols = len(table_data_1), len(table_data_1[0])
tbl_shape = slide.shapes.add_table(
    rows, cols, COL2_L, y, COL_W, Cm(rows * 1.0))
tbl = tbl_shape.table

for r, row_data in enumerate(table_data_1):
    for c, cell_text in enumerate(row_data):
        cell = tbl.cell(r, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
        run = p.add_run()
        clean = cell_text.replace("*", "")
        run.text = clean
        run.font.size = Pt(18)
        run.font.name = 'Calibri'
        is_bold = (cell_text.endswith("*") or r == 0 or
                   (r == rows - 1 and c == 0))
        run.font.bold = is_bold
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = SEC_BLUE
            run.font.color.rgb = WHITE
        elif r == rows - 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY

y += Cm(rows * 1.0) + Cm(0.8)

# Table 2 — SF-UDA
tb = add_textbox(COL2_L, y, COL_W, Cm(1.2))
tf = tb.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Table 2: SF-UDA Dice (%): BTCV CT → CHAOS MRI"
run.font.size = Pt(22)
run.font.bold = True
run.font.name = 'Calibri'
y += Cm(1.4)

table_data_2 = [
    ["Method", "Spleen", "R.Kid.", "L.Kid.", "Liver", "Avg"],
    ["Target-sup.", "94.5", "95.5", "95.3", "95.1", "95.1"],
    ["DPL [8]", "38.3", "65.3", "57.3", "63.1", "56.0"],
    ["DFG [7]", "85.1*", "92.5*", "79.6", "83.6*", "85.2*"],
    ["DiffuSAM (w/o 3D)", "76.3", "85.8", "85.2", "74.3", "80.4"],
    ["DiffuSAM", "84.6", "87.7", "86.3*", "82.2", "85.2*"],
]

rows2, cols2 = len(table_data_2), len(table_data_2[0])
tbl_shape2 = slide.shapes.add_table(
    rows2, cols2, COL2_L, y, COL_W, Cm(rows2 * 1.0))
tbl2 = tbl_shape2.table

for r, row_data in enumerate(table_data_2):
    for c, cell_text in enumerate(row_data):
        cell = tbl2.cell(r, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
        run = p.add_run()
        clean = cell_text.replace("*", "")
        run.text = clean
        run.font.size = Pt(18)
        run.font.name = 'Calibri'
        is_bold = (cell_text.endswith("*") or r == 0 or
                   (r == rows2 - 1 and c == 0))
        run.font.bold = is_bold
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = SEC_BLUE
            run.font.color.rgb = WHITE
        elif r == rows2 - 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY

y += Cm(rows2 * 1.0) + Cm(0.4)

tb = add_textbox(COL2_L, y, COL_W, Cm(1.5))
tf = tb.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Cross-slice volumetric conditioning improved average Dice by 4–5 percentage points in both settings."
run.font.size = Pt(24)
run.font.bold = True
run.font.name = 'Calibri'
y += Cm(2.0)

# t-SNE figures side by side
tsne_w = (COL_W - Cm(1)) // 2
add_image(COL2_L, y, tsne_w, 'tsne_combined_all_organs_perp3_epoch1.png')
tsne_h = add_image(COL2_L + tsne_w + Cm(1), y, tsne_w,
                   'tsne_combined_all_organs_perp3_epoch75.png')
y += tsne_h + Cm(0.2)

tb = add_textbox(COL2_L, y, COL_W, Cm(2.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = ("t-SNE of generated vs. ground-truth SAM2 memory embeddings. "
            "(A) Early training: diffuse, misaligned. "
            "(B) End of training: compact, aligned clusters.")
run.font.size = Pt(20)
run.font.italic = True
run.font.name = 'Calibri'
y += Cm(2.5)

# Inference-step figures
inf_h = add_image(COL2_L, y, COL_W, 'inference_steps_patient_36_slice_32.png')
y += inf_h + Cm(0.2)
inf_h2 = add_image(
    COL2_L, y, COL_W, 'inference_steps_patient_IMG-0007-2_slice_10.png')
y += inf_h2 + Cm(0.2)

tb = add_textbox(COL2_L, y, COL_W, Cm(2))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = ("Progressive diffusion refinement. Top: CT (BTCV). Bottom: MRI (CHAOS). "
            "Accurate segmentations emerge in only 2 diffusion steps.")
run.font.size = Pt(20)
run.font.italic = True
run.font.name = 'Calibri'
y += Cm(2.2)

# Discussion and Conclusion
dy = section_bar(COL2_L, y, COL_W, "Discussion and Conclusion")
y += dy

txt_disc = [
    ("DiffuSAM demonstrates that a compact diffusion prior over frozen "
     "foundation model features enables competitive prompt-free segmentation "
     "with minimal training resources. Current limitations include evaluation "
     "restricted to four abdominal organs across two datasets. Future work "
     "will extend to broader anatomical structures and integrate additional "
     "domain adaptation strategies to support clinical deployment.",
     False, False),
]
body_text(COL2_L, y, COL_W, txt_disc)

# ================================================================
#  FOOTER — References & Disclosures
# ================================================================
FOOTER_TOP = A0_H - Cm(9.5)

section_bar(MARGIN_L, FOOTER_TOP, BODY_W, "References & Disclosures")
ref_top = FOOTER_TOP + Cm(2.2)

refs = (
    "[1] Kirillov A et al. Segment anything. arXiv 2023; 2304.02643.\n"
    "[2] Ravi N et al. SAM 2: segment anything in images and videos. arXiv 2024; 2408.00714.\n"
    "[3] Ma J et al. Segment anything in medical images. Nat Commun 2024;15:654.\n"
    "[4] Ronneberger O et al. U-Net. MICCAI 2015:234–241.\n"
    "[5] Hatamizadeh A et al. UNETR. IEEE/CVF WACV 2022:574–584.\n"
    "[6] He X et al. FATE-SAM. arXiv 2025; 2501.09138.\n"
    "[7] Huai Z et al. DFG. IEEE TMI 2025. doi:10.1109/TMI.2025.3587733.\n"
    "[8] Chen C et al. DPL. MICCAI 2021:225–235.\n"
    "[9] Landman B et al. BTCV challenge. MICCAI Workshop 2015.\n"
    "[10] Kavur AE et al. CHAOS challenge. Med Image Anal 2021;69:101950."
)

tb = add_textbox(MARGIN_L, ref_top, COL_W + Cm(5), Cm(7))
tf = tb.text_frame
tf.word_wrap = True
for i, line in enumerate(refs.strip().split('\n')):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = line
    run.font.size = Pt(16)
    run.font.name = 'Calibri'

# Disclosures
tb = add_textbox(COL2_L, ref_top, COL_W, Cm(3))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Disclosures"
run.font.size = Pt(24)
run.font.bold = True
run.font.name = 'Calibri'
p = tf.add_paragraph()
run = p.add_run()
run.text = "The authors declare no conflicts of interest."
run.font.size = Pt(20)
run.font.name = 'Calibri'

# ================================================================
#  SAVE — versioned output under posters/
# ================================================================
POSTERS_DIR = os.path.join(PROJ, "posters")
os.makedirs(POSTERS_DIR, exist_ok=True)

BASENAME = "aimed2026_poster.pptx"


def _next_version_number(directory, basename):
    """Scan directory for ver_NNN_<basename> and return next N."""
    import re
    pattern = re.compile(r"^ver_(\d{3})_" + re.escape(basename) + r"$")
    max_n = 0
    for fname in os.listdir(directory):
        m = pattern.match(fname)
        if m:
            max_n = max(max_n, int(m.group(1)))
    return max_n + 1


ver = _next_version_number(POSTERS_DIR, BASENAME)
versioned_name = f"ver_{ver:03d}_{BASENAME}"
out = os.path.join(POSTERS_DIR, versioned_name)
prs.save(out)
print(f"Saved: {out}")
print(f"Size: {os.path.getsize(out) / 1e6:.1f} MB")
